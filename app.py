from flask import Flask, render_template, request, redirect, url_for, flash
import os
import sqlite3
import re
import logging
import zipfile
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
import cloudmersive_virus_api_client
import shutil  # Add this import at the top of your file

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['TRASH_FOLDER'] = 'trash'
app.config['ALLOWED_EXTENSIONS'] = {
    "docx", "xls", "xlsx", "ppt", "pptx", "odt", "txt", "py", "js", "c", 
    "java", "zip", "r", "bsv", "yaml", "xml", "word", "xl", "xss", "sql", "php","jpg", "png", "jpeg"
}
app.secret_key = "supersecretkey"

# Ensure directories exist
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
os.makedirs(app.config['TRASH_FOLDER'], exist_ok=True)

# Setup logging
logging.basicConfig(filename='scan.log', level=logging.INFO, format='%(asctime)s - %(message)s')

# Initialize database
def init_db():
    conn = sqlite3.connect('database.db')
    cursor = conn.cursor()
    cursor.execute("DROP TABLE IF EXISTS files;")  # Drop the table if it exists
    cursor.execute("""
        CREATE TABLE IF NOT EXISTS files (
            id INTEGER PRIMARY KEY,
            filename TEXT,
            status TEXT,
            threat_detected TEXT
        );
    """)
    conn.commit()
    conn.close()

init_db()

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in app.config['ALLOWED_EXTENSIONS']


def is_stenographed(image_path):
    """
    Check if the given image is stenographed.
    
    Parameters:
    image_path (str): The path to the image file.

    Returns:
    bool: True if the image is stenographed (malicious), False otherwise.
    """
    # File Type Verification
    file_type = subprocess.check_output(['file', image_path]).decode()
    if 'image' not in file_type:
        return False  # Not an image file

    # Metadata Analysis
    metadata = subprocess.check_output(['exiftool', image_path]).decode()
    if 'No Exif' not in metadata:
        return True  # Metadata may contain hidden data

    # Histogram Analysis
    image = cv2.imread(image_path)
    gray_image = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
    hist = cv2.calcHist([gray_image], [0], None, [256], [0, 256])
    
    # Check for unusual spikes in histogram
    if np.any(hist > 100):  # Arbitrary threshold for demonstration
        return False

    # Hex Dump Examination
    hex_dump = subprocess.check_output(['xxd', image_path]).decode()
    if 'hidden' in hex_dump:  # Example keyword to look for
        return False

    # Strings Command
    strings_output = subprocess.check_output(['strings', image_path]).decode()
    if len(strings_output) > 0:
        return False

    # If none of the checks indicate steganography
    return False

def scan_file(filepath):
    # Patterns to detect SQL injection and other malicious content
    patterns = [
        r"SELECT\s+.*\s+FROM", r"UNION\s+SELECT", r"OR\s+1=1",
        r"DROP\s+TABLE", r"INSERT\s+INTO", r"UPDATE\s+.*\s+SET",
        r"DELETE\s+FROM", r"EXEC\s+", r"SHOW\s+TABLES", r"ALTER\s+TABLE",
        r"exec\(", r"system\(", r"shell_exec\(", r"eval\(", r"base64_decode\(",
        r"assert\(", r"preg_replace\(\s*'/e'", r"cmd\s*=\s*['\"]",
        r"payload\s*=\s*['\"]*exec\(", r"username\s*:\s*['\"]*admin",  # Example for detecting admin credentials
        r"password\s*:\s*['\"]*secret",  # Example for detecting hardcoded passwords
        r"rm\s+-rf\s+/",  # Dangerous command
        r"requests\.post\("  # Detecting HTTP requests in Python
    ]

    content = ""

    try:
        # Check if the file is a ZIP file
        if filepath.endswith('.zip'):
            extract_path = os.path.join('temp_extracted', os.path.basename(filepath)[:-4])
            os.makedirs(extract_path, exist_ok=True)

            with zipfile.ZipFile(filepath, 'r') as zip_ref:
                zip_ref.extractall(extract_path)

            # Traverse each extracted file
            for root, dirs, files in os.walk(extract_path):
                for file in files:
                    file_path = os.path.join(root, file)
                    if allowed_file(file):  # Ensure the file type is allowed
                        if scan_file(file_path):  # Recursively scan the extracted file
                            logging.warning(f"Malicious file detected in ZIP: {file_path}")
                            return True

            # Clean up extracted files
            shutil.rmtree(extract_path)  # Remove the directory and its contents
            return False  # No malicious files found in the ZIP
            # Check for steganography
        elif filepath.endswith('.jpg'):
                os.remove(filepath)
                is_stenographed(filepath)
                flash(f"Warning: {file.filename} contains steganography and was rejected!")
                return redirect(url_for('index'))

        # Extract text from various file types
        elif filepath.endswith('.docx'):
            doc = Document(filepath)
            content = "\n".join([para.text for para in doc.paragraphs])

        elif filepath.endswith('.pptx'):
            ppt = Presentation(filepath)
            slide_texts = []
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text.strip())
            content = "\n".join(slide_texts)

        elif filepath.endswith('.xlsx'):
            wb = load_workbook(filepath, data_only=True)
            sheet_texts = []
            for sheet in wb.worksheets:
                for row in sheet.iter_rows():
                    for cell in row:
                        if cell.value:
                            sheet_texts.append(str(cell.value))
            content = "\n".join(sheet_texts)

        else:
            with open(filepath, 'r', errors='ignore') as f:
                content = f.read()

        # Debugging: Log extracted content
        logging.info(f"Extracted content from {filepath}: {content[:500]}")  # Logs first 500 chars

        # Scan for malicious patterns
        for pattern in patterns:
            if re.search(pattern, content, re.IGNORECASE):
                logging.warning(f"Malicious pattern detected in {filepath}")
                return True  

    except Exception as e:
        logging.error(f"Error scanning {filepath}: {e}")

    return False

def scan_metadata(filepath):
    suspicious_keywords = ["script", "eval", "base64", "exec", "shell", "system"]
    try:
        if filepath.endswith('.docx'):
            doc = Document(filepath)
            metadata = doc.core_properties
            for key, value in metadata.__dict__.items():
                if any(keyword in str(value).lower() for keyword in suspicious_keywords):
                    logging.warning(f"Suspicious metadata in {filepath}: {key} = {value}")
                    return True

        elif filepath.endswith('.pptx'):
            ppt = Presentation(filepath)
            metadata = ppt.core_properties
            for key, value in metadata.__dict__.items():
                if any(keyword in str(value).lower() for keyword in suspicious_keywords):
                    logging.warning(f"Suspicious metadata in {filepath}: {key} = {value}")
                    return True

        elif filepath.endswith('.xlsx'):
            wb = load_workbook(filepath)
            metadata = wb.properties
            for key, value in metadata.__dict__.items():
                if any(keyword in str(value).lower() for keyword in suspicious_keywords):
                    logging.warning(f"Suspicious metadata in {filepath}: {key} = {value}")
                    return True

    except Exception as e:
        logging.error(f"Error scanning metadata in {filepath}: {e}")

    return False

def heuristic_detection(filepath):
    # Heuristic-based detection for embedded objects/macros
    try:
        if filepath.endswith('.pptx'):
            ppt = Presentation(filepath)
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if shape.shape_type == 14:  # 14 is the type for embedded objects
                        logging.warning(f"Embedded object detected in {filepath}")
                        return True

        elif filepath.endswith('.xlsx'):
            wb = load_workbook(filepath, data_only=True)
            for sheet in wb.worksheets:
                for row in sheet.iter_rows():
                    for cell in row:
                        if isinstance(cell.value, str) and ("macro" in cell.value.lower() or "vba" in cell.value.lower()):
                            logging.warning(f"Potential macro detected in {filepath}")
                            return True

    except Exception as e:
        logging.error(f"Error in heuristic detection for {filepath}: {e}")

    return False

def scan_directory(directory):
    for root, _, files in os.walk(directory):
        for file in files:
            filepath = os.path.join(root, file)
            if allowed_file(file):
                threat_detected = False
                if scan_file(filepath) or scan_metadata(filepath) or heuristic_detection(filepath):
                    threat_detected = True
                    os.rename(filepath, os.path.join(app.config['TRASH_FOLDER'], file))
                    logging.warning(f"Moved {file} to trash due to suspicious content.")

                # Log the detection in the database
                conn = sqlite3.connect('database.db')
                cursor = conn.cursor()
                cursor.execute("INSERT INTO files (filename, status, threat_detected) VALUES (?, ?, ?)", 
                               (file, "Malicious" if threat_detected else "Safe", "Yes" if threat_detected else "No"))
                conn.commit()
                conn.close()

@app.route('/')
def index():
    conn = sqlite3.connect('database.db')
    cursor = conn.cursor()
    cursor.execute("SELECT * FROM files")
    files = cursor.fetchall()
    conn.close()
    return render_template('index.html', files=files)


@app.route('/upload', methods=['POST'])
def upload_file():
    if 'file' not in request.files:
        flash('No file part')
        return redirect(request.url)

    files = request.files.getlist('file')  

    for file in files:
        if file.filename == '':
            continue

        if allowed_file(file.filename):
            filepath = os.path.join(app.config['UPLOAD_FOLDER'], file.filename)
            file.save(filepath)

            if file.filename.endswith('*.zip'):
                if scan_zip_file(filepath):  # Check if the ZIP file contains malicious content
                    flash(f"Warning: {file.filename} contains malicious content and was rejected!")
                    os.remove(filepath)  # Optionally remove the ZIP file
                    return redirect(url_for('index'))

            elif scan_file(filepath):  
                os.remove(filepath)
                flash(f"Warning: {file.filename} contains suspicious content or metadata and was rejected!")
                return redirect(url_for('index'))
            else:
                flash(f'{file.filename} uploaded successfully!')

            conn = sqlite3.connect('database.db')
            cursor = conn.cursor()
            cursor.execute("INSERT INTO files (filename, status) VALUES (?, ?)", (file.filename, "Safe"))
            conn.commit()
            conn.close()
        else:
            flash(f'File type not allowed: {file.filename}')
    
    return redirect(url_for('index'))






def scan_zip_file(zip_filepath):
    """Extract and scan files from a ZIP archive for malicious content."""
    extract_path = os.path.join('temp_extracted', os.path.basename(zip_filepath)[:-4])
    os.makedirs(extract_path, exist_ok=True)

    malicious_found = False  # Flag to track if any malicious files are found

    try:
        with zipfile.ZipFile(zip_filepath, 'r') as zip_ref:
            zip_ref.extractall(extract_path)

        # Traverse each extracted file
        for root, dirs, files in os.walk(extract_path):
            for file in files:
                file_path = os.path.join(root, file)
                if allowed_file(file):  # Ensure the file type is allowed
                    if scan_file(file_path):  # Scan the file for malware
                        logging.warning(f"Malicious file detected: {file_path}")
                        os.rename(file_path, os.path.join('malicious_files', file))  # Move to a malicious folder
                        malicious_found = True  # Set the flag to True if a malicious file is found

    except Exception as e:
        logging.error(f"Error processing ZIP file: {e}")

    finally:
        # Clean up extracted files
        shutil.rmtree(extract_path)  # Remove the directory and its contents

    return malicious_found  # Return True if any malicious files were found, otherwise False



def scan_zip_file(zip_filepath):
    """Extract and scan files from a ZIP archive for malicious content."""
    extract_path = os.path.join('temp_extracted', os.path.basename(zip_filepath)[:-4])
    os.makedirs(extract_path, exist_ok=True)

    malicious_found = False  # Flag to track if any malicious files are found

    try:
        with zipfile.ZipFile(zip_filepath, 'r') as zip_ref:
            zip_ref.extractall(extract_path)

        # Traverse each extracted file
        for root, dirs, files in os.walk(extract_path):
            for file in files:
                file_path = os.path.join(root, file)
                if allowed_file(file):  # Ensure the file type is allowed
                    if scan_file(file_path):  # Scan the file for malware
                        logging.warning(f"Malicious file detected: {file_path}")
                        os.rename(file_path, os.path.join('malicious_files', file))  # Move to a malicious folder
                        malicious_found = True  # Set the flag to True if a malicious file is found

    except Exception as e:
        logging.error(f"Error processing ZIP file: {e}")

    finally:
        # Clean up extracted files
        shutil.rmtree(extract_path)  # Remove the directory and its contents

    return malicious_found  # Return True if any malicious files were found, otherwise False

def scan_file_with_api(file_path):
    # Configure API key authorization
    configuration = cloudmersive_virus_api_client.Configuration()
    configuration.api_key['Apikey'] = 'YOUR_API_KEY'
    api_instance = cloudmersive_virus_api_client.ScanApi(cloudmersive_virus_api_client.ApiClient(configuration))

    try:
        api_response = api_instance.scan_file_advanced(file_path)
        return not api_response.CleanResult  # Return True if malicious
    except Exception as e:
        print(f"Error scanning file: {e}")
        return False
@app.route('/upload-folder', methods=['POST'])
def upload_folder():
    if 'folder' not in request.files:
        flash('No folder selected')
        return redirect(request.url)

    folder_files = request.files.getlist('folder')
    if not folder_files:
        flash("No folder uploaded")
        return redirect(url_for('index'))

    folder_name = folder_files[0].filename.split('/')[0]  
    folder_path = os.path.join(app.config['UPLOAD_FOLDER'], folder_name)
    os.makedirs(folder_path, exist_ok=True)  # Ensure the main folder exists

    for file in folder_files:
        if file.filename == '':
            continue

        # Create the full path for the file
        filepath = os.path.join(folder_path, file.filename)

        # Ensure the directory for the file exists
        os.makedirs(os.path.dirname(filepath), exist_ok=True)

        # Save the file
        file.save(filepath)

        if allowed_file(file.filename) and (scan_file(filepath) or scan_metadata(filepath) or heuristic_detection(filepath)):
            os.remove(filepath)
            flash(f"Warning: {file.filename} contains malicious content or metadata and was rejected!")
        else:
            flash(f'{file.filename} uploaded successfully!')

    scan_directory(folder_path)  
    return redirect(url_for('index'))

@app.route('/delete/<int:file_id>')
def delete_file(file_id):
    conn = sqlite3.connect('database.db')
    cursor = conn.cursor()
    cursor.execute("SELECT filename FROM files WHERE id = ?", (file_id,))
    file = cursor.fetchone()

    if file:
        filename = file[0]
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        trash_path = os.path.join(app.config['TRASH_FOLDER'], filename)

        if os.path.exists(file_path):
            os.rename(file_path, trash_path)

        cursor.execute("DELETE FROM files WHERE id = ?", (file_id,))
        conn.commit()
        flash('File moved to trash!')

    conn.close()
    return redirect(url_for('index'))

if __name__ == '__main__':
    os.environ["FLASK_RUN_FROM_CLI"] = "false"  
    scan_directory(app.config['UPLOAD_FOLDER'])
    app.run(debug=False,host='0.0.0.0')