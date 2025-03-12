from pptx import Presentation

def create_malicious_pptx(filename):
    # Create a new Presentation
    ppt = Presentation()
    
    # Add a slide
    slide = ppt.slides.add_slide(ppt.slide_layouts[5])  # Title Slide layout
    
    # Add suspicious content
    suspicious_content = [
        "SELECT * FROM users WHERE username = 'admin';",
        "DROP TABLE users;",
        "<script>alert('XSS');</script>",
        "eval('some malicious code');"
    ]
    
    for line in suspicious_content:
        textbox = slide.shapes.add_textbox(left=0, top=0, width=ppt.slide_width, height=ppt.slide_height)
        text_frame = textbox.text_frame
        text_frame.text = line
    
    # Save the presentation
    ppt.save(filename)
    print(f"Malicious .pptx file '{filename}' created successfully.")

# Create a malicious .pptx file
create_malicious_pptx('malicious.pptx')
