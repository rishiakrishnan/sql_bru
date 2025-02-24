from flask import Flask, render_template, request, redirect, url_for, flash
import os
import sqlite3
import re
import logging
import zipfile
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['TRASH_FOLDER'] = 'trash'
app.config['ALLOWED_EXTENSIONS'] = {"docx", "xls", "xlsx", "ppt", "pptx", "odt", "txt", "py", "js", "c", "java", "zip", "r", "bsv", "yaml", "xml", "word", "xl", "xss", "sql", "php"}
app.secret_key = "supersecretkey"

# Ensure directories exist
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
os.makedirs(app.config['TRASH_FOLDER'], exist_ok=True)

# Setup logging
logging.basicConfig(filename='scan.log', level=logging.INFO, format='%(asctime)s - %(message)s')

# Initialize database
conn = sqlite3.connect('database.db')
cursor = conn.cursor()
cursor.execute("CREATE TABLE IF NOT EXISTS files (id INTEGER PRIMARY KEY, filename TEXT, status TEXT);")
conn.commit()
conn.close()

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in app.config['ALLOWED_EXTENSIONS']

def scan_file(filepath):
    patterns = [
        r"SELECT\s+.*\s+FROM", r"UNION\s+SELECT", r"OR\s+1=1",
        r"DROP\s+TABLE", r"INSERT\s+INTO", r"UPDATE\s+.*\s+SET",
        r"DELETE\s+FROM", r"EXEC\s+", r"SHOW\s+TABLES", r"ALTER\s+TABLE"
    ]

    content = ""

    try:
        # Extract text from .docx
        if filepath.endswith('.docx'):
            doc = Document(filepath)
            content = "\n".join([para.text for para in doc.paragraphs])

        # Extract text from .pptx
        elif filepath.endswith('.pptx'):
            ppt = Presentation(filepath)
            slide_texts = []
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text.strip())
            content = "\n".join(slide_texts)

        # Extract text from .xlsx
        elif filepath.endswith('.xlsx'):
            wb = load_workbook(filepath, data_only=True)
            sheet_texts = []
            for sheet in wb.worksheets:
                for row in sheet.iter_rows():
                    for cell in row:
                        if cell.value:
                            sheet_texts.append(str(cell.value))
            content = "\n".join(sheet_texts)

        # Extract text from .txt, .py, .js, etc.
        else:
            with open(filepath, 'r', errors='ignore') as f:
                content = f.read()

        # Debugging: Log extracted content
        logging.info(f"Extracted content from {filepath}: {content[:500]}")  # Logs first 500 chars

        # Scan for malicious patterns
        for pattern in patterns:
            if re.search(pattern, content, re.IGNORECASE):
                logging.warning(f"Malicious pattern detected in {filepath}")
                return True  

    except Exception as e:
        logging.error(f"Error scanning {filepath}: {e}")

    return False



def scan_metadata(filepath):
    suspicious_keywords = ["script", "eval", "base64", "exec", "shell", "system"]
    try:
        if filepath.endswith('.docx'):
            doc = Document(filepath)
            metadata = doc.core_properties
            for key, value in metadata.__dict__.items():
                if any(keyword in str(value).lower() for keyword in suspicious_keywords):
                    logging.warning(f"Suspicious metadata in {filepath}: {key} = {value}")
                    return True

        elif filepath.endswith('.pptx'):
            ppt = Presentation(filepath)
            metadata = ppt.core_properties
            for key, value in metadata.__dict__.items():
                if any(keyword in str(value).lower() for keyword in suspicious_keywords):
                    logging.warning(f"Suspicious metadata in {filepath}: {key} = {value}")
                    return True

        elif filepath.endswith('.xlsx'):
            wb = load_workbook(filepath)
            metadata = wb.properties
            for key, value in metadata.__dict__.items():
                if any(keyword in str(value).lower() for keyword in suspicious_keywords):
                    logging.warning(f"Suspicious metadata in {filepath}: {key} = {value}")
                    return True

    except Exception as e:
        logging.error(f"Error scanning metadata in {filepath}: {e}")

    return False

def scan_directory(directory):
    for root, _, files in os.walk(directory):
        for file in files:
            filepath = os.path.join(root, file)
            if allowed_file(file):
                if scan_file(filepath) or scan_metadata(filepath):
                    os.rename(filepath, os.path.join(app.config['TRASH_FOLDER'], file))
                    logging.warning(f"Moved {file} to trash due to suspicious content.")

@app.route('/')
def index():
    conn = sqlite3.connect('database.db')
    cursor = conn.cursor()
    cursor.execute("SELECT * FROM files")
    files = cursor.fetchall()
    conn.close()
    return render_template('index.html', files=files)

@app.route('/upload', methods=['POST'])
def upload_file():
    if 'file' not in request.files:
        flash('No file part')
        return redirect(request.url)

    files = request.files.getlist('file')  

    for file in files:
        if file.filename == '':
            continue

        if allowed_file(file.filename):
            filepath = os.path.join(app.config['UPLOAD_FOLDER'], file.filename)
            file.save(filepath)

            if file.filename.endswith('.zip'):
                extract_path = os.path.join(app.config['UPLOAD_FOLDER'], file.filename[:-4])
                os.makedirs(extract_path, exist_ok=True)

                with zipfile.ZipFile(filepath, 'r') as zip_ref:
                    zip_ref.extractall(extract_path)

                os.remove(filepath)  
                scan_directory(extract_path)  

            elif scan_file(filepath) or scan_metadata(filepath):  
                os.remove(filepath)
                flash(f"Warning: {file.filename} contains suspicious content or metadata and was rejected!")
                return redirect(url_for('index'))
            else:
                flash(f'{file.filename} uploaded successfully!')

            conn = sqlite3.connect('database.db')
            cursor = conn.cursor()
            cursor.execute("INSERT INTO files (filename, status) VALUES (?, ?)", (file.filename, "Safe"))
            conn.commit()
            conn.close()
        else:
            flash(f'File type not allowed: {file.filename}')
    
    return redirect(url_for('index'))


@app.route('/upload-folder', methods=['POST'])
def upload_folder():
    if 'folder' not in request.files:
        flash('No folder selected')
        return redirect(request.url)

    folder_files = request.files.getlist('folder')
    if not folder_files:
        flash("No folder uploaded")
        return redirect(url_for('index'))

    folder_name = folder_files[0].filename.split('/')[0]  
    folder_path = os.path.join(app.config['UPLOAD_FOLDER'], folder_name)
    os.makedirs(folder_path, exist_ok=True)  # Ensure the main folder exists

    for file in folder_files:
        if file.filename == '':
            continue

        # Create the full path for the file
        filepath = os.path.join(folder_path, file.filename)

        # Ensure the directory for the file exists
        os.makedirs(os.path.dirname(filepath), exist_ok=True)

        # Save the file
        file.save(filepath)

        if allowed_file(file.filename) and (scan_file(filepath) or scan_metadata(filepath)):
            os.remove(filepath)
            flash(f"Warning: {file.filename} contains malicious content or metadata and was rejected!")
        else:
            flash(f'{file.filename} uploaded successfully!')

    scan_directory(folder_path)  
    return redirect(url_for('index'))
@app.route('/upload-folder', methods=['POST'], endpoint='upload_folder_endpoint')
def upload_folder():
    if 'folder' not in request.files:
        flash('No folder selected')
        return redirect(request.url)

    folder_files = request.files.getlist('folder')
    if not folder_files:
        flash("No folder uploaded")
        return redirect(url_for('index'))

    folder_name = folder_files[0].filename.split('/')[0]  
    folder_path = os.path.join(app.config['UPLOAD_FOLDER'], folder_name)
    os.makedirs(folder_path, exist_ok=True)

    for file in folder_files:
        if file.filename == '':
            continue

        filepath = os.path.join(folder_path, file.filename)
        file.save(filepath)

        if allowed_file(file.filename) and (scan_file(filepath) or scan_metadata(filepath)):
            os.remove(filepath)
            flash(f"Warning: {file.filename} contains malicious content or metadata and was rejected!")
        else:
            flash(f'{file.filename} uploaded successfully!')

    scan_directory(folder_path)  
    return redirect(url_for('index'))
@app.route('/delete/<int:file_id>')
def delete_file(file_id):
    conn = sqlite3.connect('database.db')
    cursor = conn.cursor()
    cursor.execute("SELECT filename FROM files WHERE id = ?", (file_id,))
    file = cursor.fetchone()

    if file:
        filename = file[0]
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        trash_path = os.path.join(app.config['TRASH_FOLDER'], filename)

        if os.path.exists(file_path):
            os.rename(file_path, trash_path)

        cursor.execute("DELETE FROM files WHERE id = ?", (file_id,))
        conn.commit()
        flash('File moved to trash!')

    conn.close()
    return redirect(url_for('index'))
if __name__ == '__main__':
    os.environ["FLASK_RUN_FROM_CLI"] = "false"  
    scan_directory(app.config['UPLOAD_FOLDER'])
    app.run(debug=False)